
%%writefile app.py
import streamlit as st
import google.genai as genai
from google.genai import types
from PIL import Image
import pypdf
import docx
from pptx import Presentation
from gtts import gTTS
import io
import time
from datetime import datetime

st.set_page_config(
    page_title="Custom 7-Tier AI Study Suite",
    page_icon="🎓",
    layout="wide"
)

st.sidebar.title("🎓 Smart Study System")
st.sidebar.markdown("**Project Status:** 7 Core Modules Active")
st.sidebar.markdown("**Hardware:** NVIDIA GPU Accelerated")

st.sidebar.markdown("---")
st.sidebar.subheader("API Configuration")
api_key = st.sidebar.text_input("Enter Gemini API Key:", type="password", placeholder="AI Studio Key...")

if api_key:
    try:
        client = genai.Client(api_key=api_key)
        model_id = "gemini-2.5-flash"
    except Exception as e:
        st.sidebar.error("Invalid Key Format.")
        st.stop()
else:
    st.sidebar.warning("Please enter your Gemini API Key to activate the system modules.")
    st.info("💡 Tip: Get a free key from Google AI Studio using your Gmail account.")
    st.stop()

# Helper function to extract text from a Word (.docx) file
def extract_text_from_docx(file_buffer):
    doc = docx.Document(io.BytesIO(file_buffer.read()))
    return "\n".join([para.text for para in doc.paragraphs])

# Helper function to extract text from a PowerPoint (.pptx) file safely
def extract_text_from_pptx(file_buffer):
    prs = Presentation(io.BytesIO(file_buffer.read()))
    text_runs = []
    for slide in prs.slides[:5]:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_runs.append(shape.text)
    return "\n".join(text_runs)

# Resilient function to handle 503 Overload errors smoothly
def safe_stream_generator(contents_payload):
    max_retries = 3
    for attempt in range(max_retries):
        try:
            stream = client.models.generate_content_stream(model=model_id, contents=contents_payload)
            for chunk in stream:
                yield chunk.text
            return
        except Exception as e:
            if "503" in str(e) or "UNAVAILABLE" in str(e):
                if attempt < max_retries - 1:
                    time.sleep(2)
                    continue
            yield f"\n\n⚠️ **Google Server Status Notification:** The API is experiencing a high volume of traffic. Please click the button to try processing again in a moment."
            return

choice = st.sidebar.radio(
    "Select Project Module:",
    [
        "1. 🤖 AI ChatBox",
        "2. 📝 Question Auto-Solver",
        "3. 📅 AI Roadmap Generator",
        "4. 📊 PYQ Trend Analyzer",
        "5. 🎯 Handwritten Notes Mock Test",
        "6. 🎙️ Lecture Audio/Video Summarizer", # Updated Label Name
        "7. 🗣️ AI Audio Translator & Teacher"
    ]
)

# ==========================================
# MODULE 1 - AI CHATBOX
# ==========================================
if choice == "1. 🤖 AI ChatBox":
    st.header("🤖 AI ChatBox (Streaming)")
    if "messages" not in st.session_state: st.session_state.messages = []
    for msg in st.session_state.messages:
        with st.chat_message(msg["role"]): st.markdown(msg["content"])
    if prompt := st.chat_input("Ask any engineering/science doubt..."):
        with st.chat_message("user"): st.markdown(prompt)
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("assistant"):
            full_res = st.write_stream(safe_stream_generator(prompt))
            st.session_state.messages.append({"role": "assistant", "content": full_res})

# ==========================================
# MODULE 2 - QUESTION AUTO-SOLVER
# ==========================================
elif choice == "2. 📝 Question Auto-Solver":
    st.header("📝 Multimodal Question Auto-Solver")
    up_file = st.file_uploader("Upload Assignment (Image, PDF, DOCX)", type=["png", "jpg", "jpeg", "pdf", "docx"])
    if up_file:
        if "image" in up_file.type:
            img = Image.open(up_file)
            st.image(img, use_container_width=True)
            payload = [img, "Extract and solve these questions step-by-step."]
        elif "docx" in up_file.name or "officedocument" in up_file.type:
            payload = [f"Solve these questions:\n\n{extract_text_from_docx(up_file)}"]
        else:
            rdr = pypdf.PdfReader(io.BytesIO(up_file.read()))
            payload = [f"Solve these:\n\n" + "\n".join([p.extract_text() for p in rdr.pages])]
        if st.button("🚀 Auto-Solve Questions"):
            with st.chat_message("assistant"):
                st.write_stream(safe_stream_generator(payload))

# ==========================================
# MODULE 3 - AI ROADMAP GENERATOR
# ==========================================
elif choice == "3. 📅 AI Roadmap Generator":
    st.header("📅 Personalized AI Roadmap Generator")
    col1, col2 = st.columns(2)
    with col1: subj = st.text_input("Enter Subject/Topic Name:")
    with col2: e_date = st.date_input("Exam Date Target:", min_value=datetime.today())
    if st.button("🗺️ Build My Schedule"):
        days = (e_date - datetime.today().date()).days
        if days <= 0: st.error("Date must be in the future.")
        else:
            pmp = f"Create a structured step-by-step day-by-day study roadmap for: '{subj}' for {days} days."
            with st.chat_message("assistant"):
                st.write_stream(safe_stream_generator(pmp))

# ==========================================
# MODULE 4 - PYQ TREND ANALYZER
# ==========================================
elif choice == "4. 📊 PYQ Trend Analyzer":
    st.header("📊 Multi-File PYQ Pattern & Trend Analyzer")
    pyq_sub = st.text_input("Course Domain Name:")
    st.markdown("### 📥 Upload Question Papers Below (Maximum 4 Papers)")
    c1, c2 = st.columns(2)
    with c1: p1 = st.file_uploader("📄 PYQ Question Paper - 1", type=["png","jpg","jpeg","pdf","docx"], key="p1")
    with c2: p2 = st.file_uploader("📄 PYQ Question Paper - 2", type=["png","jpg","jpeg","pdf","docx"], key="p2")
    with c1: p3 = st.file_uploader("📄 PYQ Question Paper - 3", type=["png","jpg","jpeg","pdf","docx"], key="p3")
    with c2: p4 = st.file_uploader("📄 PYQ Question Paper - 4", type=["png","jpg","jpeg","pdf","docx"], key="p4")
    slots = [p for p in [p1, p2, p3, p4] if p is not None]
    if st.button("🔍 Run Cross-Paper Analytics"):
        if not pyq_sub: st.warning("Please specify the course subject name.")
        elif len(slots) == 0: st.error("Upload at least 1 question paper.")
        else:
            ctx = f"--- MASTER PYQ DATA FOR SUBJECT: {pyq_sub} ---\n"
            imgs = []
            for idx, paper in enumerate(slots):
                if "image" in paper.type:
                    imgs.append(Image.open(paper))
                    ctx += f"[Paper Slot {idx+1}: Image attachment active]\n"
                elif "docx" in paper.name or "officedocument" in paper.type:
                    ctx += f"[Paper Slot {idx+1} Word Content]:\n" + extract_text_from_docx(paper) + "\n"
                else:
                    ctx += f"[Paper Slot {idx+1} PDF Content]:\n" + "\n".join([pg.extract_text() for pg in pypdf.PdfReader(io.BytesIO(paper.read())).pages]) + "\n"
            pmp = ctx + "\nAnalyze patterns and reply with these headers: ### 🔁 1. Repeated Questions\n\n### 🔥 2. Frequently Asked Topics\n\n### 📦 3. Important Core Units\n\n### 🎯 4. Probability of Questions Appearing"
            with st.chat_message("assistant"):
                st.write_stream(safe_stream_generator(imgs + [pmp]))

# ==========================================
# MODULE 5 - HANDWRITTEN NOTES MOCK TEST (REVISED)
# ==========================================
elif choice == "5. 🎯 Handwritten Notes Mock Test":
    st.header("🎯 Notes-Based Mock Test Generator")

    # 1. Initialization
    if "quiz_stage" not in st.session_state: st.session_state.quiz_stage = "upload"

    n_file = st.file_uploader("Upload Notes (Image, PDF, or Word)", type=["png", "jpg", "jpeg", "pdf", "docx"])

    # 2. Stage: Upload & Question Generation
    if n_file and st.session_state.quiz_stage == "upload":
        if "image" in n_file.type:
            img = Image.open(n_file)
            st.image(img, width=300)
            pld = [img, "Generate 3 conceptual questions based on these notes. Return ONLY the questions. Do not include answers."]
        else:
            txt = extract_text_from_docx(n_file) if ("docx" in n_file.name or "officedocument" in n_file.type) else "\n".join([p.extract_text() for p in pypdf.PdfReader(io.BytesIO(n_file.read())).pages])
            pld = [f"Generate 3 conceptual questions from this text:\n\n{txt}\n\nReturn ONLY the questions. Do not include answers."]

        if st.button("📝 Generate Questions"):
            with st.spinner("Creating your mock test..."):
                res = client.models.generate_content(model=model_id, contents=pld)
                st.session_state.n_quiz = res.text
                st.session_state.quiz_stage = "answering"
                st.rerun()

    # 3. Stage: Answering
    if st.session_state.quiz_stage == "answering":
        st.markdown("### ✏️ Attempt the Questions:")
        st.info(st.session_state.n_quiz)
        user_answers = st.text_area("Type your answers here:", key="user_ans_area")

        if st.button("💯 Submit & Grade My Paper"):
            st.session_state.user_answers = user_answers
            st.session_state.quiz_stage = "grading"
            st.rerun()

    # 4. Stage: Grading, Correction, and Model Answer
    if st.session_state.quiz_stage == "grading":
        st.markdown("### 📊 Grading & Feedback Report")

        # Expert Teacher Prompt
        grading_prompt = f"""
        Role: You are an expert professor.
        Task: Grade the student's answers based on the original questions.

        Original Questions: {st.session_state.n_quiz}
        Student's Answers: {st.session_state.user_answers}

        Please provide:
        1. A score (out of 10).
        2. A breakdown of mistakes the student made (be specific).
        3. The correct, model answers for comparison.
        """

        with st.chat_message("assistant"):
            st.write_stream(safe_stream_generator(grading_prompt))

        if st.button("🔄 Start New Test"):
            for key in ["n_quiz", "user_answers", "quiz_stage"]:
                if key in st.session_state: del st.session_state[key]
            st.rerun()

# ==========================================
# MODULE 6 - LECTURE AUDIO/VIDEO SUMMARIZER (UPDATED FOR MULTIMEDIA)
# ==========================================
elif choice == "6. 🎙️ Lecture Audio/Video Summarizer":
    st.header("🎙️ Lecture Audio/Video-to-Notes Summarizer")
    st.caption("Upload raw recorded audio (.mp3, .wav) or video files (.mp4, .mov, .avi) to extract textbook notes natively via AI.")

    # Updated file uploader accepts both audio and video formats now
    lecture_file = st.file_uploader("Upload Lecture File (Audio or Video)", type=["mp3", "wav", "m4a", "mp4", "mov", "avi"])

    if lecture_file:
        # Dynamically render player depending on whether it's audio or video
        if "video" in lecture_file.type or lecture_file.name.endswith(('.mp4', '.mov', '.avi')):
            st.video(lecture_file)
        else:
            st.audio(lecture_file)

        if st.button("✨ Transcribe & Extract Core Notes"):
            with st.spinner("AI is carefully reading your multimedia file timeline..."):
                file_bytes = lecture_file.read()

                # Bundle the direct bytes along with the specific mime type string
                media_payload = [
                    types.Part.from_bytes(data=file_bytes, mime_type=lecture_file.type),
                    "Process this lecture recording file completely. Analyze the speech and structural concepts presented, then provide a highly detailed, textbook-style reference layout containing: 1. Core Summary 2. Important Formulas, Theorems, or Key points mentioned 3. Definitions or technical jargon to memorize."
                ]
                with st.chat_message("assistant"):
                    st.write_stream(safe_stream_generator(media_payload))

# ==========================================
# MODULE 7 - AI AUDIO TRANSLATOR & TEACHER
# ==========================================
elif choice == "7. 🗣️ AI Audio Translator & Teacher":
    st.header("🗣️ Multilingual Document Audio Teacher")
    st.caption("Upload slides or document notes and choose your preferred language. The AI will translate and teach them to you out loud.")

    lang_choice = st.selectbox(
        "Choose teaching language:",
        ["English", "Hindi (हिंदी)", "Kannada (ಕನ್ನಡ)", "French (Français)"]
    )

    lang_codes = {"English": "en", "Hindi (हिंदी)": "hi", "Kannada (ಕನ್ನಡ)": "kn", "French (Français)": "fr"}
    doc_file = st.file_uploader("Upload Study Notes (PDF, Word DOCX, or PowerPoint PPTX)", type=["pdf", "docx", "pptx"])

    if doc_file:
        if st.button("🔊 Translate & Generate Audio Lecture"):
            with st.spinner(f"Extracting notes and translating into {lang_choice}..."):
                try:
                    if "pptx" in doc_file.name:
                        extracted_notes = extract_text_from_pptx(doc_file)
                    elif "docx" in doc_file.name:
                        extracted_notes = extract_text_from_docx(doc_file)
                    else:
                        reader = pypdf.PdfReader(io.BytesIO(doc_file.read()))
                        extracted_notes = "\n".join([page.extract_text() for page in reader.pages])

                    if not extracted_notes.strip():
                        st.error("No extractable text content found.")
                    else:
                        teacher_prompt = (
                            f"You are an expert university professor teaching a class. "
                            f"Translate and explain the following document content into {lang_choice}. "
                            f"Keep your explanation concise, highly educational, easy to follow, and natural to read aloud. "
                            f"Do not include any code blocks, markdown tables, or special character patterns:\n\n{extracted_notes[:2000]}"
                        )

                        res = None
                        for attempt in range(3):
                            try:
                                res = client.models.generate_content(model=model_id, contents=teacher_prompt)
                                break
                            except Exception as api_err:
                                if attempt < 2: time.sleep(2)
                                else: raise api_err

                        translated_explanation = res.text
                        st.markdown(f"### 📝 Translated Teaching Draft ({lang_choice})")
                        st.write(translated_explanation)

                        st.markdown("---")
                        st.markdown(f"### 🔊 Spoken Audio Lecture ({lang_choice})")

                        with st.spinner("Synthesizing clear audio track..."):
                            tts = gTTS(text=translated_explanation, lang=lang_codes[lang_choice], slow=False)
                            audio_fp = io.BytesIO()
                            tts.write_to_fp(audio_fp)
                            audio_fp.seek(0)
                            st.audio(audio_fp, format="audio/mp3")
                            st.success("Audio lecture created successfully!")

                except Exception as e:
                    st.error("Google Gemini servers are highly occupied at this second. Let's wait a moment and click the button again.")



